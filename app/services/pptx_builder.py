import io
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.utils.colors import hex_to_rgbcolor
from app.services.image_search import get_image_urls_for_topic_google, download_image_to_tempfile
from flask import current_app

def apply_background(slide, theme_data):
    fill = slide.background.fill
    fill.solid()
    bg_color = theme_data.get('bg-color', '#FFFFFF')
    fill.fore_color.rgb = hex_to_rgbcolor(bg_color)

def set_font_format(text_frame, font_name, font_size, font_color, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    supported_fonts = current_app.config['SUPPORTED_FONTS']
    safe_font = font_name if font_name in supported_fonts else 'Arial'
    
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = safe_font
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = bold
            run.font.italic = italic
        paragraph.alignment = alignment
        paragraph.space_after = Pt(8)
        paragraph.space_before = Pt(0)

def create_pptx_file(slides_data, theme_data, customization):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_layout_text_only = prs.slide_layouts[1]
    content_layout_with_image = prs.slide_layouts[3]
    blank_slide_layout = prs.slide_layouts[6]
    
    temp_files_to_delete = []

    for i, slide_info in enumerate(slides_data):
        slide_type = slide_info.get('slide_type', 'content')
        title = slide_info.get('title', '')
        points = slide_info.get('points', [])
        image_query = slide_info.get('image_query', 'none')
        
        slide = None
        
        if slide_type == 'intro':
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if len(points) > 0:
                slide.placeholders[1].text = "\n".join(points)
            if slide.placeholders[1]:
                set_font_format(slide.placeholders[1].text_frame, 
                                theme_data.get('font-body', 'Arial'), 24, 
                                hex_to_rgbcolor(theme_data.get('font-color-body', '#333333')),
                                alignment=PP_ALIGN.CENTER)
        
        elif slide_type == 'thanks':
            slide = prs.slides.add_slide(blank_slide_layout)
            txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
            tf = txBox.text_frame
            tf.text = title
            set_font_format(tf, theme_data.get('font-title', 'Arial'), 44, 
                            hex_to_rgbcolor(theme_data.get('font-color-title', '#000000')), 
                            bold=True, alignment=PP_ALIGN.CENTER)
        
        else: # 'content' slide
            image_path = None
            if image_query != 'none' and customization.get('image_strategy') != 'none':
                image_urls = get_image_urls_for_topic_google(image_query)
                if image_urls:
                    for url in image_urls:
                        image_path = download_image_to_tempfile(url)
                        if image_path:
                            temp_files_to_delete.append(image_path)
                            break
                    
            if image_path:
                slide = prs.slides.add_slide(content_layout_with_image)
                slide.shapes.title.text = title
                
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                for point in points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                set_font_format(text_frame, theme_data.get('font-body', 'Calibri'), 18, 
                                hex_to_rgbcolor(theme_data.get('font-color-body', '#333333')))
                
                try:
                    placeholder = slide.placeholders[2]
                    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
                    pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                except Exception as e:
                    print(f"Error adding picture to PPT placeholder: {e}")
                    if image_path in temp_files_to_delete:
                        temp_files_to_delete.remove(image_path)
                    try:
                        os.remove(image_path)
                    except Exception as e_rm:
                        print(f"Error cleaning up failed image: {e_rm}")
            
            else:
                slide = prs.slides.add_slide(content_layout_text_only)
                slide.shapes.title.text = title
                body_shape = slide.placeholders[1]
                text_frame = body_shape.text_frame
                text_frame.clear()
                for point in points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                set_font_format(text_frame, theme_data.get('font-body', 'Calibri'), 20, 
                                hex_to_rgbcolor(theme_data.get('font-color-body', '#333333')))

        if slide:
            apply_background(slide, theme_data)
            if slide.shapes.title:
                set_font_format(slide.shapes.title.text_frame, 
                                theme_data.get('font-title', 'Arial'), 32, 
                                hex_to_rgbcolor(theme_data.get('font-color-title', '#000000')), 
                                bold=True)
                slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    
    for path in temp_files_to_delete:
        try: os.remove(path)
        except Exception as e: print(f"Error removing temp file {path}: {e}")
            
    pptx_file.seek(0)
    return pptx_file
